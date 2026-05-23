from pathlib import Path

from io import BytesIO
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from plots import build_graph_figures

# Colores
COLOR_FONDO = "F7F7F2"
COLOR_TEXTO = "1F2937"
COLOR_PRINCIPAL = "1D4E89"
COLOR_SECUNDARIO = "2A9D8F"
COLOR_ALERTA = "D95D39"
COLOR_NEUTRO = "8D99AE"

# Mejor forma de manejar direcciones de archivos
# Al hacer esto, no importa desde que directorio se ejecute el script, siempre encontrara las rutas correctas relativas al proyecto
script_dir = Path(__file__).resolve().parent
project_root = script_dir.parent

database_path = project_root / "INPUTS" / "database.sqlite"
output_path = project_root / "OUTPUTS" / "presentacion_final.pptx"

if output_path.exists():
    output_path.unlink(missing_ok=True)

def rgb(hex_color):
    hex_color = hex_color.replace("#", "")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def crear_presentacion(output_path, analyses, integrantes):
    """Crea la presentacion tomando como input el path donde se va a crear la presentacion, los analisis y los integrantes"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def nueva_slide():
        """Agrega una nueva slide vacia con fondo de color de nuestra paleta de colores"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(COLOR_FONDO)
        return slide

    def agregar_texto(slide, left, top, width, height, texto, tamano, color, negrita=False):
        """"""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = texto
        run.font.name = "Aptos"
        run.font.size = Pt(tamano)
        run.font.bold = negrita
        run.font.color.rgb = rgb(color)

    def agregar_parrafos(slide, left, top, width, height, texto, tamano=18):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True

        parrafos = [p.strip() for p in texto.split("\n\n") if p.strip()]

        for i, parrafo in enumerate(parrafos):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = parrafo
            p.space_after = Pt(10)

            for run in p.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(tamano)
                run.font.color.rgb = rgb(COLOR_TEXTO)

    # Portada
    slide = nueva_slide()
    agregar_texto(
        slide,
        Inches(0.9), Inches(1.4),
        Inches(10.5), Inches(1.0),
        "Analisis de datos para decisiones futbolisticas",
        24,
        COLOR_PRINCIPAL,
        True
    )
    agregar_texto(
        slide,
        Inches(0.9), Inches(2.4),
        Inches(10.8), Inches(0.8),
        "Scouting, perfiles de jugador, valor competitivo y apoyo a decisiones tacticas",
        15,
        COLOR_NEUTRO
    )
    agregar_texto(
        slide,
        Inches(0.9), Inches(4.3),
        Inches(10.5), Inches(1.2),
        "Presentacion con graficas e interpretaciones del analisis.",
        18,
        COLOR_TEXTO
    )

    # Analisis
    for i, analysis in enumerate(analyses, start=1):
        # Slide de grafica
        slide = nueva_slide()
        agregar_texto(
            slide,
            Inches(0.7), Inches(0.5),
            Inches(11.0), Inches(0.5),
            f"Analisis {i:02d}: {analysis['title']}",
            20,
            COLOR_PRINCIPAL,
            True
        )

        figure = analysis["figure"]

        image_stream = BytesIO()
        figure.savefig(image_stream, format="png", dpi=200, bbox_inches="tight")
        image_stream.seek(0)

        max_width = 11.0
        max_height = 5.4

        fig_width, fig_height = figure.get_size_inches()
        scale = min(max_width / fig_width, max_height / fig_height)

        final_width = fig_width * scale
        final_height = fig_height * scale

        left = 1.0 + (max_width - final_width) / 2
        top = 1.3 + (max_height - final_height) / 2

        slide.shapes.add_picture(
            image_stream,
            Inches(left), Inches(top),
            width=Inches(final_width),
            height=Inches(final_height)
        )

        # Slide de texto
        slide = nueva_slide()
        agregar_texto(
            slide,
            Inches(0.7), Inches(0.5),
            Inches(11.0), Inches(0.5),
            f"Lectura del análisis {i:02d}: {analysis['title']}",
            20,
            COLOR_SECUNDARIO,
            True
        )
        agregar_parrafos(
            slide,
            Inches(0.9), Inches(1.4),
            Inches(11.2), Inches(5.6),
            analysis["analysis_text"],
            18
        )

    # Slide final
    slide = nueva_slide()
    agregar_texto(
        slide,
        Inches(0.9), Inches(1.5),
        Inches(5.0), Inches(0.6),
        "Integrantes",
        24,
        COLOR_PRINCIPAL,
        True
    )
    agregar_parrafos(
        slide,
        Inches(1.0), Inches(2.5),
        Inches(5.0), Inches(3.0),
        "\n\n".join(integrantes),
        20
    )

    prs.save(output_path)


if __name__ == "__main__":

    grafico_01, grafico_02, grafico_03, grafico_04, grafico_05, grafico_06, grafico_07, grafico_08, grafico_09, grafico_10 = build_graph_figures(database_path) 

    analyses = [
        {
            "title": "¿Qué atributos mirar en jugadores de campo?",
            "figure": grafico_01,
            "analysis_text": """¿Qué problema intentamos resolver con nuestro análisis?

El scouting tradicional pierde tiempo evaluando decenas de atributos por jugador. Esto genera exceso de información y retrasa la toma de decisiones durante las ventanas de fichajes.

Solución propuesta:

La matriz de correlación demuestra matemáticamente que los atributos se agrupan en "bloques". Por ejemplo, Velocidad y Aceleración son prácticamente el mismo dato (0.89 de similitud), al igual que Regate y Control de Balón (0.85). No necesitamos analizar 20 variables independientes.

Implicación del negocio: 

Esto nos permite reducir a la mitad los indicadores clave de rendimiento (KPIs) Podemos crear un filtro automatizado mucho más ágil. En lugar de buscar "el jugador perfecto en 20 atributos", el equipo de scouting puede enfocarse en evaluar 4 o 5 "bloques o perfiles" fundamentales. Esto acelera el descarte de jugadores que no encajan en el club, optimiza las horas de trabajo de los visores y reduce el riesgo económico en los fichajes.
"""
        },
        {
            "title": "¿Qué atributos mirar en porteros?",
            "figure": grafico_02,
            "analysis_text": """¿Qué problema intentamos solucionar con nuestro análisis? 

La evaluación de porteros suele basarse en una sobrecarga de métricas redundantes, lo que ralentiza la identificación de talento y encarece la adquisición de datos y herramientas de scouting.

¿Qué solución se propone?

Habilidades como Estirada (Diving) y Reflejos (0.89), o Colocación (Positioning) y Manejo (Handling) (0.83), son indivisibles. En la práctica, un portero que sabe ubicarse raramente da rebotes peligrosos. El entrenamiento y el scouting deben evaluar estos pares de forma conjunta, no aislada. Se detecta una variable aislada, el salto (Jumping) nos indica que la calidad de un portero se centra más en técnica y ubicación, no en atletismo puro.

¿Qué utilidad trae al negocio el análisis?

Hacer que auditar 15 atributos distintos por portero no sea necesario da un incremento significativo en el ahorro de tiempo y dinero. El sistema de scouting se puede automatizar con solo 2 o 3 métricas clave, ahorrando cientos de horas de análisis de video y datos, asimismo ayuda a identificar “gangas” en el mercado. Al saber que el rendimiento general el predecible se pueden filtrar bases de datos en segundos para encontrar porteros de ligas menores con el mismo perfil estadístico que opciones mucho más caras.
"""
        },
        {
            "title": "¿Cuánto margen tienen los campistas jovenes?",
            "figure": grafico_03,
            "analysis_text": """¿Qué problema se trata en este análisis?

Al fichar talento joven (18 – 23 años), los clubes suelen guiarse por la intuición sobre "cuánto puede mejorar" un jugador, asumiendo riesgos financieros altos sin un parámetro de lo que es un crecimiento normal frente a uno excepcional.

¿Qué solución se propone?

Este análisis se basa exclusivamente en un conjunto de datos de ligas de primera división europea. En base al histograma, podemos observar que el margen de crecimiento (gap) de un prospecto joven incluido en este dataset tiende a concentrarse entre los 7 y 10 puntos. Este comportamiento de la muestra define nuestro estándar de mercado.

¿Qué implicaciones hay para el negocio?

El departamento de scouting debe prestar menos atención al centro de la gráfica y enfocarse exclusivamente en la "cola derecha" (prospectos con un gap mayor a 12-14 puntos). Estos son los "activos infravalorados": jugadores cuyo costo de transferencia actual es bajo porque su rendimiento inmediato no refleja su techo real. Aquí es donde el club puede comprar barato y vender caro.
"""
        },
        {
            "title": "¿Cuánto margen tienen los porteros jovenes?",
            "figure": grafico_04,
            "analysis_text": """¿Qué problema se aborda en este análisis?

Evaluar el talento y el techo de los porteros jóvenes (18-23 años) es uno de los mayores retos en el mercado, ya que su curva de maduración suele diferir de la del resto del equipo. Asumir riesgos financieros basándose solo en la intuición puede resultar en inversiones estancadas.

¿Qué solución se presenta? 

Este análisis se basa exclusivamente en un conjunto de datos de ligas de primera división europea. En base al histograma, podemos observar que el margen de crecimiento (gap) de los arqueros jóvenes incluidos en este dataset tiende a concentrarse principalmente entre los 7 y 12 puntos, con una mediana cercana a 9. Este comportamiento observado en la muestra define nuestro estándar de mercado para la posición.

¿Qué implicaciones tiene para el negocio?

El departamento de scouting debe enfocar sus recursos en la "cola derecha" de la distribución. Estos casos atípicos representan arqueros que combinan un nivel actual todavía contenido con un techo estimado mucho más alto. Identificarlos permite adquirir prospectos prometedores a un costo de transferencia inicial bajo, maximizando la plusvalía a largo plazo.
"""
        },
        {
            "title": "¿Qué distingue a los jovenes con alto potencial?",
            "figure": grafico_05,
            "analysis_text": """Los jugadores jóvenes con mayor potencial destacan principalmente por atributos técnicos y ofensivos, más que por capacidades físicas o defensivas. DRIBBLING, CURVE, BALL_CONTROL y LONG_SHOTS son las variables con mayor diferencia respecto al grupo de referencia, lo que sugiere que el principal desarrollo está más asociado con calidad técnica, creatividad y capacidad para generar jugadas ofensivas que con atributos defensivos tradicionales. También sobresalen POSITIONING y VISION, indicando que la lectura de juego y la toma de decisiones son factores importantes en perfiles con alto potencial.

Desde la perspectiva de scouting, el gráfico sugiere que varios atributos pueden agruparse en bloques similares para simplificar la evaluación de talento joven. Por ejemplo, DRIBBLING, BALL_CONTROL y CURVE reflejan capacidades técnicas relacionadas, mientras que VISION, POSITIONING y SHORT_PASSING representan comprensión táctica y participación ofensiva. Esto permitiría reducir la cantidad de métricas que un scout necesita revisar individualmente y enfocarse en atributos clave representativos de cada bloque, haciendo más eficiente la identificación de jugadores con mayor margen de desarrollo.
"""
        },
        {
            "title": "¿Qué tan cerrados son los partidos?",
            "figure": grafico_06,
            "analysis_text": """La mayor parte de los partidos del dataset se define por márgenes pequeños. Los encuentros cerrados, decididos por un solo gol, representan la categoría más frecuente con cerca del 37% del total, y al sumarlos con los empates, más del 60% de los partidos terminan con diferencias mínimas. Esto confirma que los resultados suelen ser altamente competitivos y que los partidos dominados ampliamente son mucho menos comunes. En contraste, los encuentros con diferencias de 3 o más goles representan la proporción más baja del análisis.

Desde una perspectiva táctica y de scouting, este resultado sugiere que pequeñas ventajas competitivas pueden tener un impacto decisivo en el resultado final. Factores como la toma de decisiones, la concentración, la efectividad en jugadas puntuales o la capacidad para sostener ventajas mínimas adquieren más valor que simplemente generar grandes diferencias estadísticas. Para equipos y analistas, esto implica que conviene priorizar jugadores y esquemas tácticos capaces de influir en momentos clave del partido, ya que en contextos tan cerrados los detalles terminan definiendo gran parte de los resultados.
"""
        },
        {
            "title": "¿En qué destacan las altas promesas?",
            "figure": grafico_07,
            "analysis_text": """El boxplot muestra que el grupo de Alta Promesa (Top 10%) presenta distribuciones más altas que el grupo de Resto de Jóvenes en los seis atributos evaluados. Esto se nota porque, en todos los casos, la mediana de las altas promesas queda por encima de la mediana del resto, lo que indica una ventaja consistente y no aislada en un solo atributo.

Las diferencias más marcadas aparecen en atributos como Visión, Control de Balón, Reacciones y Pase Corto. Esto sugiere que los jugadores jóvenes con mayor potencial no solo destacan por condiciones físicas, sino también por atributos técnicos y de lectura del juego. En Velocidad Sprint y Resistencia también se observa una ventaja para las altas promesas, aunque la separación entre grupos parece algo menor.

En conjunto, la gráfica sugiere que el potencial alto en jugadores jóvenes está asociado a un perfil más completo y balanceado. Es decir, las altas promesas no sobresalen solo en un atributo puntual, sino que tienden a mantener niveles superiores en varias dimensiones importantes del rendimiento."""
        },
        {
            "title": "¿A qué edad rinden mejor los jugadores?",
            "figure": grafico_08,
            "analysis_text": """De este gráfico de líneas, podemos identificar cuatro etapas que definen la carrera de un futbolista profesional:

- Fase de crecimiento: desde las edades más tempranas (16-18 años), la línea muestra una pendiente ascendente constante y pronunciada. Es en esta etapa donde los jugadores adquieren madurez táctica, experiencia y desarrollo físico.
- El "prime" o pico de rendimiento: la curva alcanza uno de sus puntos más altos y forma una "meseta" típica entre los 27 y los 31 años. Es en esta franja donde el jugador encuentra el balance perfecto: combina su máxima experiencia futbolística con un físico aún en óptimas condiciones.
- Fase de declive: a partir de los 32 años, la línea comienza un descenso evidente, reflejando la pérdida natural de capacidades físicas, como velocidad, explosividad o resistencia, que termina por afectar la valoración general.
- Outliers: El ascenso abrupto en jugadores de 35+ puede deberse a outliers. En este ambito hay que investigar más a fondo, pues puede ser que el promedio suba debido a que los jugadores más talentosos posponen su retiro, sesgando el resultado.

Si el equipo busca éxito inmediato para competir por un campeonato, debe apostar por jugadores en su "prime" (27-30 años), sabiendo que su costo será el más elevado. En cambio, si la directiva busca construir un proyecto a mediano plazo y generar plusvalía, la inversión debe ir hacia la fase de crecimiento (20-23 años), adquiriendo talento cuyo rendimiento y valor de mercado estadísticamente irán al alza."""
        },
        {
            "title": "¿Qué formaciones equilibran uso y rendimiento?",
            "figure": grafico_09,
            "analysis_text": """El bubble plot sugiere que no todas las formaciones combinan del mismo modo frecuencia de uso y rendimiento. En la zona superior derecha, donde coinciden más puntos por partido y una diferencia de gol promedio menos negativa o incluso positiva, aparecen esquemas como 4-3-3, 4-4-2 y 4-2-3-1. Entre ellas, 4-2-3-1 destaca especialmente porque combina un rendimiento competitivo con una burbuja grande, es decir, con una muestra amplia de partidos, lo que hace más confiable la señal.

También aparecen algunas formaciones con buen desempeño, pero con menor respaldo de muestra, como 3-5-2 y otras variantes cercanas al bloque superior. Esto sugiere que pueden ser opciones interesantes, pero conviene ser más prudente al interpretarlas, ya que una menor frecuencia de uso hace más difícil separar un patrón estable de un resultado puntual.

En el extremo opuesto, 5-3-2, 4-5-1 y 4-1-4-1 se ubican en la zona de peor balance, con menos puntos por partido y una diferencia de gol promedio más negativa. En conjunto, la gráfica sugiere que no basta con mirar qué formación se usa más: conviene priorizar esquemas que mantengan un equilibrio entre frecuencia y resultados, y en ese sentido 4-2-3-1 y 4-4-2 parecen alternativas más sólidas que otras formaciones menos consistentes."""
        },
        {
            "title": "¿Qué formaciones son más seguras o más riesgosas?",
            "figure": grafico_10,
            "analysis_text": """Desde una perspectiva táctica, el gráfico sugiere que las formaciones pueden clasificarse según su perfil de riesgo: esquemas agresivos para contextos donde ganar es prioritario, y esquemas más sólidos para torneos largos o partidos donde minimizar derrotas es más importante. Tanto para entrenadores como analistas, esto permite evaluar formaciones no solo por porcentaje de victorias, sino por la distribución completa de resultados. Así, una formación con menos derrotas y más empates puede ser más valiosa en ciertos contextos competitivos que otra con más victorias, pero también con mayor inestabilidad en sus resultados.

Este análisis es una interpretación más “superficial” ya que la cantidad de formaciones analizadas no se toma en cuenta. Esto sesga los resultados y hace que algunas formaciones (como la 3-2-3-2) parezcan aparentemente más efectivas, pero esto se debe a que su muestra es más reducida. A diferencia de la diapositiva anterior donde se consideran muestras más amplias y con mejor representación, por lo que sus conclusiones pueden ser más definitivas y tener un mejor enfoque de comparación.
"""
        }
    ]

    integrantes = [
        "Ossian Ramirez",
        "Estefania Elizabeth",
        "Ian Pascual",
        "Karim Jacob"
    ]

    crear_presentacion(
        output_path=output_path,
        analyses=analyses,
        integrantes=integrantes
    )